import subprocess
import requests

from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation

from src.constants import PPTX_TEMPLATE_PATH, DOWNLOAD_DIR


def download_image(url: str, save_dir: Path, filename: str) -> str:
    if not url:
        return ""
    save_dir.mkdir(parents=True, exist_ok=True)
    local_path = save_dir / filename
    try:
        resp = requests.get(url, timeout=10)
        resp.raise_for_status()
        with open(local_path, "wb") as f:
            f.write(resp.content)
        return str(local_path)
    except Exception as e:
        print(f"Failed to download {url}: {e}")
        return ""


def replace_images_by_alt_text(pptx_path: str, alt_text_to_image_path: dict[str, str], output_pptx_path: str):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in list(slide.shapes):
            alt_text = shape.name
            if not alt_text:
                continue

            key = alt_text[1:-1].strip() if alt_text.startswith("{") and alt_text.endswith("}") else alt_text

            if key in alt_text_to_image_path and alt_text_to_image_path[key]:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                # noinspection PyProtectedMember
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(alt_text_to_image_path[key], left, top, width, height)

    prs.save(output_pptx_path)


def replace_text_in_pptx(pptx_path: str, replacements: dict[str, str], output_path: str):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            placeholder = f"{{{key}}}"
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, value)
    prs.save(output_path)


def pptx_to_pdf(input_pptx: str, output_dir: str):
    subprocess.run([
        "libreoffice", "--headless", "--convert-to", "pdf",
        input_pptx, "--outdir", output_dir
    ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)


def pdf_to_images(pdf_path: str, output_dir: str, dpi: int = 200):
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    images = convert_from_path(pdf_path, dpi=dpi)
    image_paths = []
    for i, image in enumerate(images, start=1):
        image_path = output_dir / f"slide_{i}.png"
        image.save(image_path, 'PNG')
        image_paths.append(str(image_path))
    return image_paths


def fill_powerpoint_template(template_data: dict[str, str]):
    alt_text_to_image_path = {
        key: path for key, path in template_data.items()
        if key.endswith("_image") or "_image_" in key
    }

    pptx_path = template_data['temp_pptx_path']
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    dir_path = str(Path(pptx_path).parent)

    replace_images_by_alt_text(str(PPTX_TEMPLATE_PATH), alt_text_to_image_path, pptx_path)
    replace_text_in_pptx(pptx_path, template_data, pptx_path)
    pptx_to_pdf(pptx_path, dir_path)
    return pdf_to_images(pdf_path, str(DOWNLOAD_DIR))
