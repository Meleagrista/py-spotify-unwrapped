import os
import subprocess
import tempfile
from collections import Counter
from datetime import datetime
from pathlib import Path

import requests
import spotipy
from dotenv import load_dotenv
from pdf2image import convert_from_path
from pptx import Presentation
from spotipy.oauth2 import SpotifyOAuth

from src.constants import TEMPLATE_DIR

load_dotenv()

YOUR_APP_CLIENT_ID = os.getenv("YOUR_APP_CLIENT_ID")
YOUR_APP_CLIENT_SECRET = os.getenv("YOUR_APP_CLIENT_SECRET")

MAX_LEN = 15

def download_image(url: str, save_dir: Path, filename: str) -> str:
    if not url:
        return ""
    save_dir.mkdir(parents=True, exist_ok=True)
    local_path = save_dir / filename
    try:
        resp = requests.get(url, timeout=10)
        resp.raise_for_status()
        with open(local_path, "wb") as f:
            f.write(resp.content)
        return str(local_path)
    except Exception as e:
        print(f"Failed to download {url}: {e}")
        return ""

def replace_images_by_alt_text(pptx_path: str, alt_text_to_image_path: dict[str, str], output_pptx_path: str):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in list(slide.shapes):
                alt_text = shape.name
                if not alt_text:
                    continue

                if alt_text.startswith("{") and alt_text.endswith("}"):
                    key = alt_text[1:-1].strip()
                else:
                    key = alt_text

                if key in alt_text_to_image_path and alt_text_to_image_path[key]:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    # noinspection PyProtectedMember
                    sp = shape._element
                    sp.getparent().remove(sp)
                    slide.shapes.add_picture(alt_text_to_image_path[key], left, top, width, height)

    prs.save(output_pptx_path)

def replace_text_in_pptx(pptx_path: str, replacements: dict[str, str], output_path: str):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            placeholder = f"{{{key}}}"
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, value)
    prs.save(output_path)

def pptx_to_pdf(input_pptx: str, output_dir: str):
    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        input_pptx,
        "--outdir", output_dir
    ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

def pdf_to_images(pdf_path: str, output_dir: str, dpi: int = 200):
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    images = convert_from_path(pdf_path, dpi=dpi)
    image_paths = []
    for i, image in enumerate(images, start=1):
        image_path = output_dir / f"slide_{i}.png"
        image.save(image_path, 'PNG')
        image_paths.append(str(image_path))
    return image_paths

def current_user_top_genre(sp: spotipy.Spotify, user_data: dict[str, str], top_k: int = 5):
    all_genres = []
    total_artists = 0
    count = 1
    while True:
        offset = 50 * count
        response = sp.current_user_top_artists(limit=50, offset=offset, time_range='short_term')
        artists = response.get('items', [])
        if len(artists) == 0:
            break
        total_artists += len(artists)
        for artist in artists:
            all_genres.extend(artist.get('genres', []))
        count += 1
    genre_counts = Counter(all_genres)
    top_genres = genre_counts.most_common(top_k)
    for idx, genre in enumerate(top_genres, start=1):
        genre_key = f"genre_{idx}"
        user_data[genre_key] = genre[0].capitalize()
    return user_data

def main():
    sp = spotipy.Spotify(auth_manager=SpotifyOAuth(
        client_id=YOUR_APP_CLIENT_ID,
        client_secret=YOUR_APP_CLIENT_SECRET,
        redirect_uri="http://127.0.0.1:8000/callback",
        scope="user-library-read,user-top-read,user-read-recently-played",
    ))

    slides_dir = TEMPLATE_DIR / "slides"
    template_path = TEMPLATE_DIR / "spotify-wrapped-template.pptx"

    user_data = {"month": datetime.now().strftime("%B").upper()}

    def truncate(text: str, max_length: int) -> str:
        return text[:max_length].strip() + "..." if len(text) > max_length else text

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)

        images_dir = tmpdir
        pptx_path = tmpdir / "spotify-wrapped-temp-output.pptx"
        pdf_path = tmpdir / "spotify-wrapped-temp-output.pdf"

        top_tracks = sp.current_user_top_tracks(limit=5, time_range='short_term')
        for idx, track in enumerate(top_tracks['items'], start=1):
            track_name_key = f"track_{idx}"
            track_artist_key = f"track_artist_{idx}"
            track_image_key = f"track_image_{idx}"
            track_name = truncate(track['name'], MAX_LEN)
            track_artists = truncate(", ".join(artist['name'] for artist in track['artists']), MAX_LEN)
            album_images = track.get('album', {}).get('images', [])
            track_image_url = album_images[0]['url'] if album_images else ""
            local_track_image = download_image(track_image_url, images_dir, f"track_{idx}.jpg")
            user_data[track_name_key] = track_name
            user_data[track_artist_key] = track_artists
            user_data[track_image_key] = local_track_image

        top_artists = sp.current_user_top_artists(limit=5, time_range='short_term')
        for idx, artist in enumerate(top_artists['items'], start=1):
            artist_key = f"artist_{idx}"
            artist_image_key = f"artist_image_{idx}"
            artist_name = truncate(artist['name'], MAX_LEN)
            artist_images = artist.get('images', [])
            artist_image_url = artist_images[0]['url'] if artist_images else ""
            local_artist_image = download_image(artist_image_url, images_dir, f"artist_{idx}.jpg")
            user_data[artist_key] = artist_name
            user_data[artist_image_key] = local_artist_image

        user_data = current_user_top_genre(sp, user_data, 5)

        alt_text_to_image_path = {
            key: path for key, path in user_data.items()
            if key.endswith("_image") or "_image_" in key
        }

        replace_images_by_alt_text(str(template_path), alt_text_to_image_path, str(pptx_path))
        replace_text_in_pptx(str(pptx_path), user_data, str(pptx_path))
        pptx_to_pdf(str(pptx_path), str(tmpdir))
        pdf_to_images(str(pdf_path), str(slides_dir))

if __name__ == "__main__":
    main()
